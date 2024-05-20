#!/usr/bin/env python
import logging
import docx2txt
from zipfile import ZipFile
from subprocess import Popen, PIPE, STDOUT
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import StringIO
from pptx import Presentation


class Converter:

    def __init__(self, verbose=False):
        try:
            self.logger = logging.getLogger("SimplyEmail.Converter")
            self.verbose = verbose
        except Exception as e:
            print(e)

    def convert_docx_to_txt(self, path):
        """
        A very simple conversion function which returns unicode text for parsing.

        path = The path to the file
        """
        try:
            text = docx2txt.process(path)
            self.logger.debug(f"Converted docx to text: {path}")
            return text
        except Exception as e:
            self.logger.error(f"Failed to convert DOCX to text: {e}")
            return ""

    def convert_doc_to_txt(self, path):
        """
        A very simple conversion function which returns text for parsing.

        path = The path to the file
        """
        try:
            cmd = ['antiword', path]
            p = Popen(cmd, stdout=PIPE, stderr=STDOUT)
            stdout, stderr = p.communicate()
            return stdout.decode('ascii', 'ignore')
        except Exception as e:
            self.logger.error(f"Failed to convert DOC to text: {e}")
            return ""

    def convert_pptx_to_txt(self, path):
        """
        Converts PPTX to text.

        path = The path to the file
        """
        prs = Presentation(path)
        text_runs = ""
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs += str(run.text) + ' '
            return text_runs
        except Exception as e:
            self.logger.error(f"Failed to convert pptx: {e}")
            return text_runs if text_runs else ""

    def convert_pdf_to_txt(self, path):
        """
        A very simple conversion function which returns text for parsing from PDF.

        path = The path to the file
        """
        try:
            rsrcmgr = PDFResourceManager()
            retstr = StringIO()
            laparams = LAParams()
            device = TextConverter(rsrcmgr, retstr, laparams=laparams)
            fp = open(path, 'rb')
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            for page in PDFPage.get_pages(fp):
                interpreter.process_page(page)
            text = retstr.getvalue()
            fp.close()
            device.close()
            retstr.close()
            return text
        except Exception as e:
            self.logger.error(f"Failed to convert PDF to text: {e}")
            return ""

    def convert_xlsx_to_csv(self, path):
        """
        Converts XLSX to CSV using the xlsx2csv tool.

        path = The path to the file
        """
        self.logger.debug(f"convert_xlsx_to_csv on file: {path}")
        try:
            cmd = ['xlsx2csv', path]
            p = Popen(cmd, stdout=PIPE, stderr=STDOUT)
            stdout, stderr = p.communicate()
            text = stdout.decode('ascii', 'ignore')
            return text
        except Exception as e:
            self.logger.error(f"Failed to convert XLSX to CSV: {e}")
            return ""

    def convert_zip_to_text(self, path, rawtext=True):
        """
        Converts ZIP file contents to text.

        path = The path to the file
        rawtext = Whether to extract raw text from files
        """
        try:
            self.logger.debug(f"Attempting to unzip file: {path}")
            input_zip = ZipFile(path)
            if rawtext:
                text = ""
                for name in input_zip.namelist():
                    try:
                        text += input_zip.read(name).decode('utf-8')
                    except Exception as e:
                        print(e)
                self.logger.debug(f"Unzip of file completed (raw text): {path}")
                return text
            else:
                return {name: input_zip.read(name) for name in input_zip.namelist()}
        except Exception as e:
            self.logger.error(f"Failed to unzip file: {e}")
            return ""
